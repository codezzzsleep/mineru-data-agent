from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptInches


OUTPUT_DIR = Path("examples/office_files")


def build_docx() -> None:
    doc = Document()
    doc.add_heading("Industry Standard Review Packet", level=1)
    doc.add_paragraph("Document ID: DOCX-2026-STD-02")
    doc.add_paragraph("Effective Date: 2026-05-23")
    doc.add_paragraph(
        "Purpose: evaluate Word document structure extraction, clause tables, risk fields, "
        "recommendations, and document-level provenance."
    )
    doc.add_heading("1. Compliance Matrix", level=2)
    table = doc.add_table(rows=1, cols=4)
    header = table.rows[0].cells
    header[0].text = "Clause"
    header[1].text = "Requirement"
    header[2].text = "Evidence"
    header[3].text = "Risk"
    rows = [
        ("1.1", "Every run must keep trace logs.", "trace.json", "low"),
        ("1.2", "Structured outputs must include tables and key-values.", "result.json", "medium"),
        ("1.3", "Warnings must not be hidden from reviewers.", "quality.issues", "medium"),
    ]
    for row in rows:
        cells = table.add_row().cells
        for index, value in enumerate(row):
            cells[index].text = value
    doc.add_heading("2. Exception Handling", level=2)
    doc.add_paragraph("Risk: missing provenance should trigger review instead of silent pass.")
    doc.add_paragraph("Recommendation: inspect recovery_decision before exporting downstream chunks.")
    doc.add_paragraph("Owner: Data Governance Office")
    doc.save(OUTPUT_DIR / "industry_standard_review.docx")


def build_pptx() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Workflow Agent Review"
    left = PptInches(0.7)
    top = PptInches(1.3)
    width = PptInches(8.2)
    height = PptInches(1.1)
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.text = "Inspection Date: 2026-05-23"
    text_frame.add_paragraph().text = "Goal: verify slide-level provenance, execution matrices, and recovery notes."

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Execution Matrix"
    table_shape = slide2.shapes.add_table(4, 4, PptInches(0.6), PptInches(1.3), PptInches(8.4), PptInches(2.0))
    table = table_shape.table
    headers = ["Step", "Tool", "Output", "Recovery"]
    for col, header in enumerate(headers):
        table.cell(0, col).text = header
    rows = [
        ["1", "Native PPTX extractor", "slide text", "keep slide index"],
        ["2", "Validator", "quality report", "flag warnings"],
        ["3", "Retrieval exporter", "jsonl chunks", "skip noise"],
    ]
    for row_index, row in enumerate(rows, start=1):
        for col, value in enumerate(row):
            table.cell(row_index, col).text = value

    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Findings"
    findings = slide3.shapes.add_textbox(PptInches(0.7), PptInches(1.3), PptInches(8.2), PptInches(2.5))
    findings.text_frame.text = "Anomaly: slide decks often lose section order during plain text extraction."
    findings.text_frame.add_paragraph().text = "Recommendation: preserve slide index and table rows in content blocks."
    findings.text_frame.add_paragraph().text = "Risk: charts still require a visual model follow-up."
    prs.save(OUTPUT_DIR / "workflow_agent_review.pptx")


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    build_docx()
    build_pptx()
    (OUTPUT_DIR / "README.md").write_text(
        "# Office Fixtures\n\nSynthetic DOCX/PPTX fixtures for file-level Data Agent evidence.\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
